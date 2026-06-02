"""
Deterministic text extraction for PDF, DOCX, and PPTX files.
Returns a list of text-unit dicts with page/slide/section metadata.
"""

import logging
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text(file_path: str, file_type: str) -> list[dict]:
    """Dispatch extraction based on file type."""
    file_type = file_type.lower().lstrip(".")
    if file_type == "pdf":
        return _extract_pdf(file_path)
    if file_type == "docx":
        return _extract_docx(file_path)
    if file_type == "pptx":
        return _extract_pptx(file_path)
    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(file_path: str) -> list[dict]:
    units = []
    with fitz.open(file_path) as doc:
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                units.append({
                    "text": text,
                    "page_number": page_index,
                    "slide_number": None,
                    "section_title": None,
                    "source_type": "pdf_page",
                })
    logger.info("Extracted %d pages from PDF %s", len(units), file_path)
    return units


def _extract_docx(file_path: str) -> list[dict]:
    doc = Document(file_path)
    units = []
    current_heading = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name.lower().startswith("heading"):
            current_heading = text
        units.append({
            "text": text,
            "page_number": None,
            "slide_number": None,
            "section_title": current_heading,
            "source_type": "docx_paragraph",
        })

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                units.append({
                    "text": row_text,
                    "page_number": None,
                    "slide_number": None,
                    "section_title": current_heading,
                    "source_type": "docx_table_row",
                })

    logger.info("Extracted %d units from DOCX %s", len(units), file_path)
    return units


def _extract_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    units = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        slide_title = None
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            # Use the slide title placeholder if available
            if (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format
                and shape.placeholder_format.idx == 0
            ):
                # ✓ Use centralized config for text preview length
                from app.config import settings
                slide_title = shape.text.strip()[:settings.text_preview_chars]
        if texts:
            from app.config import settings
            units.append({
                "text": "\n".join(texts),
                "page_number": None,
                "slide_number": idx,
                "section_title": slide_title or (texts[0][:settings.text_preview_chars] if texts else None),
                "source_type": "pptx_slide",
            })

    logger.info("Extracted %d slides from PPTX %s", len(units), file_path)
    return units
